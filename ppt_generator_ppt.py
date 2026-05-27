"""
PowerPoint generator - creates presentations from structured content
"""

from pptx import Presentation
from pptx.util import Inches
import templates_ppt as templates
import config_ppt as config


class PPTGenerator:
    """Generate PowerPoint presentations from structured content"""

    def __init__(self, output_path: str = 'output.pptx'):
        self.prs = Presentation()
        self.prs.slide_width = config.SLIDE_DIMENSIONS['width']
        self.prs.slide_height = config.SLIDE_DIMENSIONS['height']
        self.output_path = output_path
        self.blank_slide_layout = self.prs.slide_layouts[6]  # Blank layout

    def add_title_slide(self, title: str, subtitle: str = ''):
        """Add title slide"""
        slide = templates.TitleSlide(self.prs, self.blank_slide_layout)
        slide.add_title(title)
        if subtitle:
            slide.add_content(subtitle)

    def add_content_slide(self, title: str, content: str):
        """Add content slide with title and body text"""
        slide = templates.ContentSlide(self.prs, self.blank_slide_layout)
        slide.add_title(title)
        slide.add_content(content)

    def add_bullet_slide(self, title: str, bullets: list):
        """Add bullet point slide"""
        slide = templates.BulletSlide(self.prs, self.blank_slide_layout)
        slide.add_title(title)
        slide.add_content(bullets)

    def add_slide_from_content(self, slide_data: dict):
        """Add slide based on content type"""
        title = slide_data.get('title', 'Slide')
        content = slide_data.get('content', '')
        slide_type = slide_data.get('type', 'content')

        if slide_type == 'title':
            self.add_title_slide(title, content if isinstance(content, str) else '')
        elif slide_type == 'bullets' and isinstance(content, list):
            self.add_bullet_slide(title, content)
        else:
            self.add_content_slide(title, content if isinstance(content, str) else str(content))

    def save(self, filename: str = None):
        """Save presentation to file"""
        output_file = filename or self.output_path
        self.prs.save(output_file)
        return output_file


def generate_from_slides(slides: list, output_path: str = 'output.pptx') -> str:
    """Convenience function to generate PPT from slides list"""
    generator = PPTGenerator(output_path)
    for slide in slides:
        generator.add_slide_from_content(slide)
    return generator.save()
