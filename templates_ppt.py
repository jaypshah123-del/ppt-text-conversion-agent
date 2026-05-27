"""
Slide templates for different content types
"""

from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import config_ppt as config


class SlideTemplate:
    """Base class for slide templates"""

    def __init__(self, prs, slide_layout):
        self.prs = prs
        self.slide = prs.slides.add_slide(slide_layout)
        self.shapes = self.slide.shapes

    def add_title(self, title: str):
        """Add title to slide"""
        raise NotImplementedError

    def add_content(self, content):
        """Add content to slide"""
        raise NotImplementedError


class TitleSlide(SlideTemplate):
    """Title slide template"""

    def add_title(self, title: str):
        """Add centered title to slide"""
        left = config.MARGINS['left']
        top = Inches(2.5)
        width = config.SLIDE_DIMENSIONS['width'] - config.MARGINS['left'] - config.MARGINS['right']
        height = Inches(2)

        text_box = self.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = config.FONTS['title']['size']
        p.font.name = config.FONTS['title']['name']
        p.font.bold = config.FONTS['title']['bold']
        p.font.color.rgb = config.FONTS['title']['color']
        p.alignment = PP_ALIGN.CENTER

    def add_content(self, content):
        """Add subtitle content"""
        left = config.MARGINS['left']
        top = Inches(4.5)
        width = config.SLIDE_DIMENSIONS['width'] - config.MARGINS['left'] - config.MARGINS['right']
        height = Inches(2)

        text_box = self.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = str(content)
        p.font.size = config.FONTS['subtitle']['size']
        p.font.name = config.FONTS['subtitle']['name']
        p.font.color.rgb = config.FONTS['subtitle']['color']
        p.alignment = PP_ALIGN.CENTER


class ContentSlide(SlideTemplate):
    """Content slide with title and body text"""

    def add_title(self, title: str):
        """Add title to top of slide"""
        left = config.MARGINS['left']
        top = config.MARGINS['top']
        width = config.SLIDE_DIMENSIONS['width'] - config.MARGINS['left'] - config.MARGINS['right']
        height = Inches(1)

        text_box = self.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = config.FONTS['heading']['size']
        p.font.name = config.FONTS['heading']['name']
        p.font.bold = config.FONTS['heading']['bold']
        p.font.color.rgb = config.FONTS['heading']['color']

    def add_content(self, content):
        """Add body text content"""
        left = config.MARGINS['left']
        top = Inches(1.5)
        width = config.SLIDE_DIMENSIONS['width'] - config.MARGINS['left'] - config.MARGINS['right']
        height = Inches(5.5)

        text_box = self.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

        p = text_frame.paragraphs[0]
        p.text = str(content)
        p.font.size = config.FONTS['body']['size']
        p.font.name = config.FONTS['body']['name']
        p.font.color.rgb = config.FONTS['body']['color']
        p.level = 0


class BulletSlide(SlideTemplate):
    """Slide with title and bullet points"""

    def add_title(self, title: str):
        """Add title to top of slide"""
        left = config.MARGINS['left']
        top = config.MARGINS['top']
        width = config.SLIDE_DIMENSIONS['width'] - config.MARGINS['left'] - config.MARGINS['right']
        height = Inches(1)

        text_box = self.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = config.FONTS['heading']['size']
        p.font.name = config.FONTS['heading']['name']
        p.font.bold = config.FONTS['heading']['bold']
        p.font.color.rgb = config.FONTS['heading']['color']

    def add_content(self, bullets: list):
        """Add bullet points"""
        left = config.MARGINS['left']
        top = Inches(1.5)
        width = config.SLIDE_DIMENSIONS['width'] - config.MARGINS['left'] - config.MARGINS['right']
        height = Inches(5.5)

        text_box = self.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.TOP

        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = bullet
            p.level = 0
            p.font.size = config.FONTS['bullet']['size']
            p.font.name = config.FONTS['bullet']['name']
            p.font.color.rgb = config.FONTS['bullet']['color']
